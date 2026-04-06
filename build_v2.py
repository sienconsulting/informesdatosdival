"""
Informe DIVAL 2025 - Versió 2 ampliada (~36 slides)
Construït amb python-pptx des de la plantilla base.
"""

import json, sys, copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

sys.stdout.reconfigure(encoding='utf-8', errors='replace')

BASE     = r'C:\Users\Pere\OneDrive - Sien Consulting\PLANIFICACION\00.NUEVA ORGANIZACION\01. Proyectos\Diputación Valencia\2024_SMART OFFICE\Informe datos'
TEMPLATE = os.path.join(BASE, 'Base', 'Informe_Trimestral_DIVAL_Mock.pptx')
METRICS  = os.path.join(BASE, 'metrics_dival.json')
OUTPUT   = os.path.join(BASE, 'Informe_DIVAL_2025.pptx')

TEAL   = RGBColor(0x00, 0x3B, 0x42)
CYAN   = RGBColor(0x04, 0xA2, 0xB6)
GREEN  = RGBColor(0x95, 0xC2, 0x1E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF4, 0xF4)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
MGRAY  = RGBColor(0xBB, 0xBB, 0xBB)
ORANGE = RGBColor(0xE8, 0x72, 0x22)

MESOS = ['Gen','Feb','Mar','Abr','Mai','Jun','Jul','Ago','Set','Oct','Nov','Des']

with open(METRICS, encoding='utf-8') as f:
    M = json.load(f)

print("Mètriques carregades. Construint informe...")

# ── Helpers ────────────────────────────────────────────────────────────────────
def fmt(n, dec=0):
    try:
        n = float(n)
        if dec == 0: return f"{int(round(n)):,}".replace(",",".")
        return f"{n:,.{dec}f}".replace(",","X").replace(".",",").replace("X",".")
    except: return str(n)

def fmtp(p, dec=1):
    if p is None: return "N/D"
    try:
        p = float(p)
        return f"{'+'if p>0 else ''}{p:.{dec}f}%"
    except: return str(p)

def monthly_list(key):
    """Return 12-element list from a dict keyed '1'..'12'."""
    d = M[key]
    return [d[str(i)] for i in range(1, 13)]

prs = Presentation(TEMPLATE)
W, H = prs.slide_width, prs.slide_height
N_TPL = len(prs.slides)
blank = prs.slide_layouts[6]

# ── Slide building helpers ─────────────────────────────────────────────────────
def slide_copy(src_idx):
    tpl = prs.slides[src_idx]
    ns = prs.slides.add_slide(blank)
    sp = ns.shapes._spTree
    for c in list(sp): sp.remove(c)
    for c in tpl.shapes._spTree: sp.append(copy.deepcopy(c))
    return ns

def add_rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, sz=12, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tx

def tb_lines(slide, lines, l, t, w, h, sz=10, color=DGRAY):
    """Multi-paragraph textbox. lines = list of str or (str, bold, size)."""
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        txt, bld, s = (item, False, sz) if isinstance(item, str) else \
                      (item[0], item[1], item[2] if len(item)>2 else sz)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = txt
        r.font.size = Pt(s); r.font.bold = bld; r.font.color.rgb = color
    return tx

def chart_add(slide, ctype, cats, series, l, t, w, h,
              legend=True, lpos='bottom'):
    cd = ChartData()
    cd.categories = cats
    for nm, vals in series: cd.add_series(nm, vals)
    cf = slide.shapes.add_chart(ctype, l, t, w, h, cd)
    ch = cf.chart
    ch.has_legend = legend and len(series) > 1
    if ch.has_legend:
        from pptx.enum.chart import XL_LEGEND_POSITION as LP
        ch.legend.position = {'bottom': LP.BOTTOM, 'right': LP.RIGHT,
                              'top': LP.TOP}.get(lpos, LP.BOTTOM)
        ch.legend.include_in_layout = False
    pal = [TEAL, CYAN, GREEN, ORANGE, MGRAY]
    for i, s in enumerate(ch.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = pal[i % len(pal)]
    return ch

def line_chart(slide, cats, series, l, t, w, h, legend=True):
    cd = ChartData()
    cd.categories = cats
    for nm, vals in series: cd.add_series(nm, vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, l, t, w, h, cd)
    ch = cf.chart
    ch.has_legend = legend
    if ch.has_legend:
        from pptx.enum.chart import XL_LEGEND_POSITION as LP
        ch.legend.position = LP.BOTTOM; ch.legend.include_in_layout = False
    pal = [TEAL, CYAN, GREEN, ORANGE]
    for i, s in enumerate(ch.series):
        s.format.line.color.rgb = pal[i % len(pal)]
        s.format.line.width = Pt(2)
    return ch

def kpi(slide, l, t, w, h, val, lbl, bg=TEAL):
    add_rect(slide, l, t, w, h, fill=bg)
    tb(slide, val, l, t+Cm(0.4), w, Cm(1.8), sz=26, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, lbl, l, t+Cm(2.0), w, Cm(1.2), sz=9, bold=False,
       color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

def header(slide, sec_n, title, subtitle=""):
    add_rect(slide, Cm(0), Cm(0), W, Cm(2.0), fill=TEAL)
    tb(slide, title.upper(), Cm(0.5), Cm(0.1), W-Cm(4), Cm(1.4), sz=14, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, Cm(0.5), Cm(1.4), W-Cm(4), Cm(0.7), sz=10, color=CYAN)
    add_rect(slide, W-Cm(3.2), Cm(0.2), Cm(1.3), Cm(1.3), fill=CYAN)
    tb(slide, str(sec_n), W-Cm(3.2), Cm(0.15), Cm(1.3), Cm(1.3),
       sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def footer(slide):
    add_rect(slide, Cm(0), H-Cm(0.55), W, Cm(0.55), fill=TEAL)
    tb(slide, "Font: INE, GVA Turisme, Seguretat Social · Smart Office Diputació de València",
       Cm(0.3), H-Cm(0.55), W-Cm(0.6), Cm(0.55), sz=7, color=WHITE)

def divider(slide, num, title, subtitle, color_accent=CYAN):
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs: r.text = ''
    tb(slide, num, Cm(1.5), Cm(2.3), Cm(3), Cm(2.5), sz=80, bold=True, color=color_accent)
    tb(slide, title, Cm(1.5), Cm(4.8), Cm(18), Cm(2.2), sz=34, bold=True, color=WHITE)
    tb(slide, subtitle, Cm(1.5), Cm(7.0), Cm(22), Cm(0.8), sz=10.5, color=color_accent)

def conclusions_slide(slide, sec_label, items, accent=CYAN):
    add_rect(slide, Cm(0), Cm(0), W, H, fill=TEAL)
    tb(slide, sec_label, Cm(1), Cm(0.3), W-Cm(2), Cm(1.1), sz=13, bold=True, color=accent)
    tb(slide, "Conclusions", Cm(1), Cm(1.3), W-Cm(2), Cm(1.1), sz=22, bold=True, color=WHITE)
    y = Cm(2.8)
    for title, body in items:
        add_rect(slide, Cm(0.8), y+Cm(0.1), Cm(0.45), Cm(0.45), fill=accent)
        tb(slide, title, Cm(1.6), y, W-Cm(2.5), Cm(0.65), sz=11, bold=True, color=WHITE)
        tb(slide, body,  Cm(1.6), y+Cm(0.6), W-Cm(2.5), Cm(0.85), sz=9.5,
           color=RGBColor(0xCC,0xCC,0xCC), wrap=True)
        y += Cm(1.6)

# ── Pre-compute data ───────────────────────────────────────────────────────────
years = ['2019','2020','2021','2022','2023','2024','2025']
rec_hist = M['historico_receptor']
int_hist = M['historico_interno']
rec_vals = [rec_hist[y] for y in years]
int_vals = [int_hist[y] for y in years]
total_hist = [r+n for r,n in zip(rec_vals, int_vals)]

rec_m25 = monthly_list('receptor_mensual_2025')
rec_m24 = monthly_list('receptor_mensual_2024')
int_m25 = monthly_list('interno_mensual_2025')
int_m24 = monthly_list('interno_mensual_2024')
p25     = monthly_list('pernoctaciones_mensual_2025')
p24     = monthly_list('pernoctaciones_mensual_2024')
oc25    = monthly_list('ocupacion_hoteles_mensual_2025')
oc24    = monthly_list('ocupacion_hoteles_mensual_2024')
em_m    = monthly_list('estancia_media_mensual')

ss      = M['afiliados_ss']  # dict keyed by date string
ss_dates = ['2025-01-01','2025-04-01','2025-07-01','2025-10-01']
quarters = ['Q1 (Gen)','Q2 (Abr)','Q3 (Jul)','Q4 (Oct)']
aloj_v  = [ss[d]['Alojamiento'] for d in ss_dates]
rest_v  = [ss[d]['Servicios de comidas y bebidas'] for d in ss_dates]
agenc_v = [ss[d]['Agencias de viaje'] for d in ss_dates]

comarca_d = {k: v for k,v in M['turistas_por_comarca_2025'].items()}
comarca_sorted = sorted(comarca_d.items(), key=lambda x: x[1]['total'], reverse=True)

hotel_cat = {k if k.strip() else 'SENSE CATEG.': v
             for k,v in M['hoteles_por_categoria'].items()}
hotel_com = sorted(M['hoteles_por_comarca'].items(), key=lambda x: x[1]['plazas'], reverse=True)
vut_mun   = sorted(M['vut_por_municipio'].items(), key=lambda x: x[1]['count'], reverse=True)
vut_com   = sorted(M['vut_por_comarca'].items(),   key=lambda x: x[1]['count'], reverse=True)
camp_mun  = sorted(M.get('campings_por_municipio',{}).items(), key=lambda x: x[1]['plazas'], reverse=True)
rural_mun = sorted(M.get('casasrurales_por_municipio',{}).items(), key=lambda x: x[1]['count'], reverse=True)

top10i = M['top10_paises_2025']
top10n = M['top10_origenes_nacionales_2025']

total_plazas = (M['hoteles_plazas'] + M['vut_plazas'] +
                M['campings_plazas'] + M['casasrurales_plazas'] + M['albergues_plazas'])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
print("S1: Portada")
s = slide_copy(0)
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs: r.text = ''
tb(s, "DIPUTACIÓ DE VALÈNCIA", Cm(1), Cm(4.2), Cm(20), Cm(2.0),
   sz=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
tb(s, "Informe de dades turístiques · ENE–DIC 2025", Cm(1), Cm(6.2), Cm(20), Cm(1.2),
   sz=17, color=CYAN, align=PP_ALIGN.LEFT)
tb(s, "Smart Office · Diputació de València",
   Cm(1), Cm(7.1), Cm(20), Cm(0.8), sz=10, color=MGRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ÍNDEX
# ══════════════════════════════════════════════════════════════════════════════
print("S2: Índex")
s = slide_copy(1)
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs: r.text = ''
tb(s, "ÍNDEX", Cm(1), Cm(0.5), Cm(15), Cm(1.4), sz=28, bold=True, color=TEAL)
secs = [("01","Afluència turística","Evolució 2019-2025, mercats, comarques"),
        ("02","Ocupació","Pernoctacions, estada, grau d'ocupació"),
        ("03","Oferta turística","Hotels, VUT, càmpings, cases rurals"),
        ("04","Treball turístic","Afiliats SS per sector i estacionalitat"),
        ("05","Conclusions","Anàlisi i perspectiva 2026")]
bw, bh, gap = Cm(4.5), Cm(5.0), Cm(0.35)
for i,(num,tit,sub) in enumerate(secs):
    lx = Cm(0.4) + i*(bw+gap)
    add_rect(s, lx, Cm(2.0), bw, bh, fill=TEAL if i%2==0 else CYAN)
    tb(s, num, lx, Cm(2.1), bw, Cm(1.5), sz=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, tit, lx, Cm(3.5), bw, Cm(1.0), sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
    tb(s, sub, lx, Cm(4.5), bw, Cm(1.8), sz=8.5, color=RGBColor(0xDD,0xDD,0xDD), align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESUM EXECUTIU
# ══════════════════════════════════════════════════════════════════════════════
print("S3: Resum executiu")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=LGRAY)
add_rect(s, Cm(0), Cm(0), W, Cm(1.8), fill=TEAL)
tb(s, "RESUM EXECUTIU · ENE–DIC 2025", Cm(0.5), Cm(0.2), W-Cm(1), Cm(1.4),
   sz=15, bold=True, color=WHITE)

kpi_data = [
    (fmt(M['total_turistas_2025']),    f"Turistes totals\n{fmtp(M['chg_total'])} vs 2024",       TEAL),
    (fmt(M['total_receptor_2025']),    f"Internacionals (rècord)\n{fmtp(M['chg_receptor'])}",     CYAN),
    (fmt(M['total_interno_2025']),     f"Nacionals\n{fmtp(M['chg_interno'])} vs 2024",            GREEN),
    (fmt(M['pernoctaciones_2025']),    f"Pernoctacions\n{fmtp(M['chg_pernoctaciones'])} vs 2024", TEAL),
    (f"{M['ocupacion_hoteles_2025_avg']:.1f}%", f"Ocupació hotelera\n{fmtp(M['chg_ocupacion'])} vs 2024", CYAN),
    (fmt(total_plazas),               f"Places registrades\n(hotels+VUT+càmpings+rural)",        GREEN),
]
bw2, bh2, gp2 = Cm(7.6), Cm(2.35), Cm(0.3)
for i,(v,l,bg) in enumerate(kpi_data):
    lx = Cm(0.4) + (i%3)*(bw2+gp2)
    ty = Cm(2.0) + (i//3)*(bh2+gp2)
    kpi(s, lx, ty, bw2, bh2, v, l, bg=bg)

total = M['total_turistas_2025']
summary = (
    f"La província de València va registrar {fmt(total)} turistes en 2025 ({fmtp(M['chg_total'])} respecte a 2024). "
    f"El turisme internacional va assolir un rècord históric: {fmt(M['total_receptor_2025'])} visitants ({fmtp(M['chg_receptor'])}). "
    f"Les pernoctacions hoteleres van sumar {fmt(M['pernoctaciones_2025'])} amb una estada mitjana de "
    f"{fmt(M['estancia_media_total'],2)} nits. L'oferta reglada supera les {fmt(total_plazas)} places."
)
tb(s, summary, Cm(0.5), H-Cm(2.0), W-Cm(1), Cm(1.7), sz=9.5, color=DGRAY, wrap=True)
footer(s)

# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓ 1 — AFLUÈNCIA TURÍSTICA
# ══════════════════════════════════════════════════════════════════════════════
print("S4: Divisor Afluència")
s = slide_copy(2)
divider(s, "01", "AFLUÈNCIA\nTURÍSTICA",
        f"Total 2025: {fmt(M['total_turistas_2025'])} turistes  ·  Internacional: {fmtp(M['chg_receptor'])}  ·  Nacional: {fmtp(M['chg_interno'])}")

print("S5: Históric total")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Evolució histórica total 2019–2025 (nacional + internacional)")
ch = chart_add(s, XL_CHART_TYPE.COLUMN_STACKED, years,
               [('Nacional', int_vals), ('Internacional', rec_vals)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8), legend=True, lpos='bottom')
ch.series[0].format.fill.fore_color.rgb = CYAN
ch.series[1].format.fill.fore_color.rgb = TEAL
tb(s, f"Total 2025:\n{fmt(total_hist[-1])}\nturistes", Cm(16.3), Cm(3.5), Cm(6.0), Cm(2.5),
   sz=15, bold=True, color=TEAL)
tb(s, f"Des de 2019 el total s'ha doblat: de {fmt(total_hist[0])} a {fmt(total_hist[-1])} visitants "
      f"(+{fmt((total_hist[-1]/total_hist[0]-1)*100,1)}%). El turisme internacional ha crescut "
      f"de forma ininterrompuda, mentre el nacional assoleix el seu segon màxim históric en 2023 "
      f"i es corregeix en 2025. En termes de pes, el turisme internacional ja representa el "
      f"{fmt(rec_vals[-1]/total_hist[-1]*100,1)}% del total.",
   Cm(16.3), Cm(6.2), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S6: Internacional históric")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Turisme receptor (internacional) — sèrie histórica")
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, years,
               [('Turistes internacionals', rec_vals)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8), legend=False)
ch.series[0].format.fill.fore_color.rgb = TEAL
kpi(s, Cm(16.3), Cm(2.5), Cm(6.0), Cm(2.8), fmt(M['total_receptor_2025']),
    f"Turistes internacionals 2025\n{fmtp(M['chg_receptor'])} vs 2024", bg=TEAL)
tb(s, f"El turisme internacional va assolir {fmt(M['total_receptor_2025'])} visitants en 2025, "
      f"un màxim históric i un creixement del {fmtp(M['chg_receptor'])} respecte a 2024. "
      f"La taxa de creixement anual composta (CAGR) des de 2019 és del "
      f"{fmt((rec_vals[-1]/rec_vals[0])**(1/6)*100-100,1)}%. "
      f"França ({fmt(top10i[0]['turistas'])}), Itàlia ({fmt(top10i[1]['turistas'])}, {fmtp(top10i[1]['chg'])}) "
      f"i Regne Unit ({fmt(top10i[2]['turistas'])}) lideren els mercats emissors. "
      f"EUA (+{fmtp(top10i[6]['chg'])}) i Polònia (+{fmtp(top10i[7]['chg'])}) destaquen per crescement.",
   Cm(16.3), Cm(5.5), Cm(6.0), Cm(7.0), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S7: Nacional históric")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Turisme intern (nacional) — sèrie histórica")
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, years,
               [('Turistes nacionals', int_vals)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8), legend=False)
ch.series[0].format.fill.fore_color.rgb = CYAN
kpi(s, Cm(16.3), Cm(2.5), Cm(6.0), Cm(2.8), fmt(M['total_interno_2025']),
    f"Turistes nacionals 2025\n{fmtp(M['chg_interno'])} vs 2024", bg=CYAN)
tb(s, f"El turisme nacional cap a València va registrar {fmt(M['total_interno_2025'])} visitants en 2025 "
      f"({fmtp(M['chg_interno'])} vs 2024). Madrid ({fmt(top10n[0]['turistas'])}, {fmtp(top10n[0]['chg'])}) "
      f"continua liderant, seguit d'Alacant ({fmt(top10n[1]['turistas'])}) i Castelló "
      f"({fmt(top10n[2]['turistas'])}). La demanda nacional representa el "
      f"{fmt(M['total_interno_2025']/M['total_turistas_2025']*100,1)}% del total. "
      f"Cal contextualitzar: el nivell de 2025 supera en valor absolut el de 2021 i 2022.",
   Cm(16.3), Cm(5.5), Cm(6.0), Cm(7.0), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S8: Distribució mensual 2025")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Distribució mensual 2025 — Nacional vs Internacional")
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, MESOS,
               [('Nacional', int_m25), ('Internacional', rec_m25)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8), legend=True, lpos='bottom')
ch.series[0].format.fill.fore_color.rgb = CYAN
ch.series[1].format.fill.fore_color.rgb = TEAL
peak_mes = MESOS[rec_m25.index(max(rec_m25))]
tb(s, f"L'anàlisi mensual revela estacionalitat estiuenca clara. "
      f"El turisme internacional supera el nacional d'abril a desembre. "
      f"El màxim internacional s'assoleix al {peak_mes} ({fmt(max(rec_m25))} visitants). "
      f"La tardor mostra bon comportament: l'octubre registra {fmt(rec_m25[9])} turistes internacionals, "
      f"un dels millors mesos de l'any per a aquest segment.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S9: Comparativa mensual internacional 2024 vs 2025")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Turisme internacional: comparativa mensual 2024 vs 2025")
ch = line_chart(s, MESOS, [('2024', rec_m24), ('2025', rec_m25)],
                Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8))
ch.series[0].format.line.color.rgb = MGRAY
ch.series[1].format.line.color.rgb = TEAL
best_m = max(range(12), key=lambda i: rec_m25[i]-rec_m24[i])
tb(s, f"El turisme internacional creix en tots els mesos de 2025 respecte a 2024. "
      f"El major increment es produeix al {MESOS[best_m]} (+{fmt(rec_m25[best_m]-rec_m24[best_m])} turistes). "
      f"L'any acumula {fmt(sum(rec_m25))} visitants internacionals vs {fmt(sum(rec_m24))} de 2024 "
      f"({fmtp((sum(rec_m25)/sum(rec_m24)-1)*100)}). "
      f"La tendència alcista és especialment marcada a tardor-hivern.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S10: Comparativa mensual nacional 2024 vs 2025")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Turisme nacional: comparativa mensual 2024 vs 2025")
ch = line_chart(s, MESOS, [('2024', int_m24), ('2025', int_m25)],
                Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8))
ch.series[0].format.line.color.rgb = MGRAY
ch.series[1].format.line.color.rgb = CYAN
worst_m = min(range(12), key=lambda i: int_m25[i]-int_m24[i])
tb(s, f"El turisme nacional de 2025 presenta variacions negatives en la majoria de mesos. "
      f"La caiguda és especialment pronunciada al {MESOS[worst_m]} "
      f"({fmtp((int_m25[worst_m]/int_m24[worst_m]-1)*100)}). "
      f"Cal recordar que 2024 va registrar màxims excepcionals. "
      f"El nivell de 2025 ({fmt(sum(int_m25))} visitants) supera els de 2021 i 2022, "
      f"situant-se com el tercer millor any de la sèrie histórica.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S11: Top 10 mercats internacionals")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Top 10 mercats emissors internacionals 2025")
i_names = [x['pais'] for x in top10i]
i_vals  = [x['turistas'] for x in top10i]
i_chgs  = [x['chg'] for x in top10i]
ch = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, i_names, [('Turistes 2025', i_vals)],
               Cm(0.5), Cm(2.3), Cm(14.0), Cm(11.0), legend=False)
ch.series[0].format.fill.fore_color.rgb = TEAL
rows = [f"{n}: {fmt(v)} ({fmtp(c)})" for n,v,c in zip(i_names,i_vals,i_chgs)]
tb_lines(s, rows, Cm(14.8), Cm(2.5), Cm(7.5), Cm(11.0), sz=9.5, color=DGRAY)
footer(s)

print("S12: Top 10 mercats nacionals")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Top 10 mercats emissors nacionals 2025")
n_names = [x['origen'] for x in top10n]
n_vals  = [x['turistas'] for x in top10n]
n_chgs  = [x['chg'] for x in top10n]
ch = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, n_names, [('Turistes nacionals 2025', n_vals)],
               Cm(0.5), Cm(2.3), Cm(14.0), Cm(11.0), legend=False)
ch.series[0].format.fill.fore_color.rgb = CYAN
rows = [f"{n}: {fmt(v)} ({fmtp(c)})" for n,v,c in zip(n_names,n_vals,n_chgs)]
tb_lines(s, rows, Cm(14.8), Cm(2.5), Cm(7.5), Cm(11.0), sz=9.5, color=DGRAY)
footer(s)

print("S13: Desglose per comarques — total")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Desglose per comarques — Turistes totals 2025")
top12c = comarca_sorted[:12]
c_names = [c[0][:22] for c in top12c]
c_int   = [c[1]['receptor'] for c in top12c]
c_nac   = [c[1]['interno']  for c in top12c]
ch = chart_add(s, XL_CHART_TYPE.BAR_STACKED, c_names,
               [('Nacional', c_nac), ('Internacional', c_int)],
               Cm(0.5), Cm(2.3), Cm(16.5), Cm(11.0), legend=True, lpos='bottom')
ch.series[0].format.fill.fore_color.rgb = CYAN
ch.series[1].format.fill.fore_color.rgb = TEAL
top1, top2 = comarca_sorted[0], comarca_sorted[1]
tb(s, f"La comarca de {top1[0]} concentra {fmt(top1[1]['total'])} turistes "
      f"({fmt(top1[1]['pct_internacional'],1)}% internacionals). "
      f"La Safor ({fmt(top2[1]['total'])}) és la principal destinació costanera. "
      f"Les comarques de l'Horta aporten volum significatiu tant de turisme nacional com internacional.",
   Cm(17.2), Cm(3.0), Cm(5.3), Cm(7.0), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S14: Perfil internacional per comarca")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "01", "Afluència turística", "Percentatge de turisme internacional per comarca 2025")
com_pct = sorted([(c[0][:22], c[1]['pct_internacional'])
                  for c in comarca_sorted if c[1]['total'] > 5000],
                 key=lambda x: x[1], reverse=True)
cp_names = [x[0] for x in com_pct]
cp_vals  = [x[1] for x in com_pct]
ch = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, cp_names, [('% Internacional', cp_vals)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(11.0), legend=False)
ch.series[0].format.fill.fore_color.rgb = TEAL
tb(s, f"{cp_names[0]} és la comarca amb major internacionalització ({fmt(cp_vals[0],1)}%). "
      f"Les comarques de l'interior com Los Serranos i La Ribera Alta presenten perfils "
      f"sorprenentment internacionals. {cp_names[-1]} ({fmt(cp_vals[-1],1)}%) és la que "
      f"té el perfil més domèstic, reflectint el turisme familiar de La Safor.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S15: Conclusions afluència")
s = prs.slides.add_slide(blank)
conclusions_slide(s, "01  AFLUÈNCIA TURÍSTICA", [
    ("Rècord históric de turisme internacional",
     f"La província assoleix {fmt(M['total_receptor_2025'])} visitants internacionals "
     f"(+{fmt(M['chg_receptor'],1)}%), el màxim de la sèrie histórica. França, Itàlia i Regne Unit lideren."),
    ("Diversificació dels mercats emissors",
     f"Itàlia (+{fmtp(top10i[1]['chg'])}), EUA (+{fmtp(top10i[6]['chg'])}) i Polònia (+{fmtp(top10i[7]['chg'])}) "
     f"creixen amb força, ampliant la base internacional i reduint dependències."),
    ("Turisme nacional en correcció moderada",
     f"El turisme intern registra {fmt(M['total_interno_2025'])} visitants ({fmtp(M['chg_interno'])} vs 2024), "
     f"tot i que supera els nivells de 2021 i 2022. Madrid domina amb {fmt(top10n[0]['turistas'])}."),
    ("Comarca de València centralitza l'afluència",
     f"Concentra el {fmt(comarca_sorted[0][1]['total']/M['total_turistas_2025']*100,1)}% del total provincial, "
     f"seguida de La Safor ({fmt(comarca_sorted[1][1]['total'])}) com a principal destí costaner."),
    ("Menor estacionalitat gràcies al turisme internacional",
     f"El flux internacional es distribueix millor al llarg de l'any, "
     f"amb octubre assolint {fmt(rec_m25[9])} turistes internacionals."),
])

# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓ 2 — OCUPACIÓ
# ══════════════════════════════════════════════════════════════════════════════
print("S16: Divisor Ocupació")
s = slide_copy(2)
divider(s, "02", "OCUPACIÓ",
        f"{fmt(M['pernoctaciones_2025'])} pernoctacions  ·  Estada {fmt(M['estancia_media_total'],2)} nits  ·  Ocupació {M['ocupacion_hoteles_2025_avg']:.1f}%")

print("S17: KPIs ocupació")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=LGRAY)
add_rect(s, Cm(0), Cm(0), W, Cm(1.8), fill=TEAL)
tb(s, "02  OCUPACIÓ · Principals indicadors 2025", Cm(0.5), Cm(0.2), W-Cm(1), Cm(1.4),
   sz=14, bold=True, color=WHITE)
ocu_kpis = [
    (fmt(M['pernoctaciones_2025']),   f"Pernoctacions hoteleres\n{fmtp(M['chg_pernoctaciones'])} vs 2024", TEAL),
    (fmt(M['viajeros_2025']),         f"Viajants als hotels\n{fmtp((M['viajeros_2025']/M['viajeros_2024']-1)*100)} vs 2024", CYAN),
    (f"{M['estancia_media_total']:.2f}", "Estada mitjana total\n(nits per visitant)", GREEN),
    (f"{M['estancia_media_nacional']:.2f}",     "Estada turista nacional\n(nits)", TEAL),
    (f"{M['estancia_media_internacional']:.2f}", "Estada turista internacional\n(nits)", CYAN),
    (f"{M['ocupacion_hoteles_2025_avg']:.1f}%",  f"Ocupació hotelera anual\n{fmtp(M['chg_ocupacion'])} vs 2024", GREEN),
]
for i,(v,l,bg) in enumerate(ocu_kpis):
    lx = Cm(0.4) + (i%3)*(bw2+gp2)
    ty = Cm(2.0) + (i//3)*(bh2+gp2)
    kpi(s, lx, ty, bw2, bh2, v, l, bg=bg)
tb(s, f"Agost: màxim pernoctacions {fmt(max(p25))}  ·  Màxima ocupació {max(oc25):.1f}% (agost)  ·  Mínim ocupació {min(oc25):.1f}% (gener)",
   Cm(0.5), H-Cm(0.9), W-Cm(1), Cm(0.75), sz=9.5, color=DGRAY)
footer(s)

print("S18: Pernoctaciones 2024 vs 2025")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "02", "Ocupació", "Pernoctacions hoteleres mensuals — comparativa 2024 vs 2025")
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, MESOS,
               [('2024', p24), ('2025', p25)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8), legend=True, lpos='bottom')
ch.series[0].format.fill.fore_color.rgb = MGRAY
ch.series[1].format.fill.fore_color.rgb = TEAL
kpi(s, Cm(16.3), Cm(2.5), Cm(6.0), Cm(2.8), fmt(sum(p25)),
    f"Total pernoctacions 2025\n{fmtp((sum(p25)/sum(p24)-1)*100)} vs 2024", bg=TEAL)
tb(s, f"Les pernoctacions hoteleres de 2025 totalitzen {fmt(sum(p25))}, "
      f"{fmtp((sum(p25)/sum(p24)-1)*100)} respecte als {fmt(sum(p24))} de 2024. "
      f"Agost és el mes de major activitat amb {fmt(max(p25))} pernoctacions. "
      f"Les temporades de primavera i tardor mostren comportaments similars entre anys.",
   Cm(16.3), Cm(5.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S19: Ocupació hotelera 2024 vs 2025")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "02", "Ocupació", "Grau d'ocupació hotelera mensual — comparativa 2024 vs 2025")
ch = line_chart(s, MESOS, [('2024', oc24), ('2025', oc25)],
                Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8))
ch.series[0].format.line.color.rgb = MGRAY
ch.series[1].format.line.color.rgb = TEAL
kpi(s, Cm(16.3), Cm(2.5), Cm(6.0), Cm(2.8),
    f"{M['ocupacion_hoteles_2025_avg']:.1f}%",
    f"Ocupació hotelera 2025\n{fmtp(M['chg_ocupacion'])} vs {M['ocupacion_hoteles_2024_avg']:.1f}% 2024",
    bg=TEAL)
tb(s, f"L'ocupació hotelera de 2025 es va situar en el {M['ocupacion_hoteles_2025_avg']:.1f}% "
      f"de mitjana anual ({fmtp(M['chg_ocupacion'])} vs {M['ocupacion_hoteles_2024_avg']:.1f}% de 2024). "
      f"El pic de l'agost arriba al {max(oc25):.1f}%, mentre que gener és el mínim ({min(oc25):.1f}%). "
      f"La tardor manté bons nivells: octubre al {oc25[9]:.1f}%.",
   Cm(16.3), Cm(5.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S20: Estada mitjana mensual")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "02", "Ocupació", "Estada mitjana mensual i per procedència")
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, MESOS,
               [('Estada mitjana (nits)', em_m)],
               Cm(0.5), Cm(2.3), Cm(12.0), Cm(10.8), legend=False)
ch.series[0].format.fill.fore_color.rgb = TEAL
kpi(s, Cm(12.8), Cm(2.5), Cm(4.8), Cm(2.6), f"{M['estancia_media_nacional']:.2f}",
    "Nits estada\nturista nacional", bg=CYAN)
kpi(s, Cm(12.8), Cm(5.3), Cm(4.8), Cm(2.6), f"{M['estancia_media_internacional']:.2f}",
    "Nits estada\nturista internacional", bg=TEAL)
kpi(s, Cm(12.8), Cm(8.1), Cm(4.8), Cm(2.6), f"{M['estancia_media_total']:.2f}",
    "Estada mitjana\ntotal", bg=GREEN)
tb(s, f"L'estada total és de {M['estancia_media_total']:.2f} nits. "
      f"El turista nacional pernocta de mitjana {M['estancia_media_nacional']:.2f} nits, "
      f"enfront de les {M['estancia_media_internacional']:.2f} del turista estranger. "
      f"El màxim mensual s'assoleix a l'agost ({max(em_m):.2f} nits).",
   Cm(12.8), H-Cm(1.6), Cm(9.7), Cm(1.3), sz=9, color=DGRAY, wrap=True)
footer(s)

print("S21: Conclusions ocupació")
s = prs.slides.add_slide(blank)
conclusions_slide(s, "02  OCUPACIÓ", [
    ("Pernoctacions estables al voltant dels 17,5 milions",
     f"Les {fmt(M['pernoctaciones_2025'])} pernoctacions de 2025 ({fmtp(M['chg_pernoctaciones'])} vs 2024) "
     f"demostren una demanda hotelera resilient."),
    ("Ocupació hotelera sòlida al 58,5%",
     f"L'ocupació del {M['ocupacion_hoteles_2025_avg']:.1f}% mostra la competitivitat de la planta hotelera, "
     f"amb pic d'agost del {max(oc25):.1f}% i octubre per damunt del {oc25[9]:.1f}%."),
    ("Diferencial estada nacional/internacional",
     f"El turista nacional pernocta {M['estancia_media_nacional']:.2f} nits vs {M['estancia_media_internacional']:.2f} "
     f"del turista estranger, reflectint dos models de viatge molt diferenciats."),
    ("Apartaments turístics complementen l'oferta",
     f"L'ocupació d'apartaments turístics al {M['ocupacion_apartamentos_2025_avg']:.1f}% "
     f"indica la rellevància del segment extrahoteler en el mix d'allotjament."),
])

# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓ 3 — OFERTA TURÍSTICA
# ══════════════════════════════════════════════════════════════════════════════
print("S22: Divisor Oferta")
s = slide_copy(2)
divider(s, "03", "OFERTA\nTURÍSTICA",
        f"{fmt(total_plazas)} places registrades  ·  {M['hoteles_total']} hotels  ·  {M['vut_total']} VUT  ·  {M['campings_total']} càmpings  ·  {M['casasrurales_total']} cases rurals",
        color_accent=GREEN)

print("S23: Resum oferta")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=LGRAY)
add_rect(s, Cm(0), Cm(0), W, Cm(1.8), fill=TEAL)
tb(s, "03  OFERTA TURÍSTICA · Resum provincial 2025", Cm(0.5), Cm(0.2), W-Cm(1), Cm(1.4),
   sz=14, bold=True, color=WHITE)
of_kpis = [
    (fmt(M['hoteles_total']),       f"Hotels i allotjaments\n{fmt(M['hoteles_plazas'])} places", TEAL),
    (fmt(M['vut_total']),           f"Vivendes d'ús turístic\n{fmt(M['vut_plazas'])} places",   CYAN),
    (fmt(M['campings_total']),      f"Càmpings\n{fmt(M['campings_plazas'])} places",             GREEN),
    (fmt(M['casasrurales_total']),  f"Cases rurals\n{fmt(M['casasrurales_plazas'])} places",     TEAL),
    (fmt(M['albergues_total']),     f"Albergs\n{fmt(M['albergues_plazas'])} places",             CYAN),
    (fmt(total_plazas),            "Total places reglades\nprovíncia de València",              GREEN),
]
for i,(v,l,bg) in enumerate(of_kpis):
    lx = Cm(0.4) + (i%3)*(bw2+gp2)
    ty = Cm(2.0) + (i//3)*(bh2+gp2)
    kpi(s, lx, ty, bw2, bh2, v, l, bg=bg)
tb(s, f"Les VUT representen el {fmt(M['vut_plazas']/total_plazas*100,1)}% de les places totals. "
      f"La Safor, l'Horta i la comarca de València concentren la major part de l'oferta provincial.",
   Cm(0.5), H-Cm(0.9), W-Cm(1), Cm(0.75), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S24: Hotels per categoria")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "03", "Oferta turística", "Establiments hotelers per categoria — Establiments i places")
cat_order = ['1 ESTRELLA','2 ESTRELLAS','3 ESTRELLAS','4 ESTRELLAS','5 ESTRELLAS','PENSIÓN','SENSE CATEG.']
cat_items = [(k,v) for k,v in hotel_cat.items()]
cat_items.sort(key=lambda x: cat_order.index(x[0]) if x[0] in cat_order else 99)
c_labels = [k.replace(' ESTRELLAS','★').replace(' ESTRELLA','★') for k,_ in cat_items]
c_estab  = [v['establecimientos'] for _,v in cat_items]
c_plazas = [v['plazas'] for _,v in cat_items]
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, c_labels,
               [('Establiments', c_estab), ('Places (÷10)', [p//10 for p in c_plazas])],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8), legend=True, lpos='bottom')
ch.series[0].format.fill.fore_color.rgb = TEAL
ch.series[1].format.fill.fore_color.rgb = CYAN
stars4 = hotel_cat.get('4 ESTRELLAS', {})
tb(s, f"Els {M['hoteles_total']} establiments hotelers sumen {fmt(M['hoteles_plazas'])} places. "
      f"La categoria 4 estrelles ({stars4.get('establecimientos',0)} estab., {fmt(stars4.get('plazas',0))} places) "
      f"concentra el major volum de places ({fmt(stars4.get('plazas',0)/M['hoteles_plazas']*100,1)}% del total). "
      f"Les pensions i allotjaments d'1 estrella representen l'oferta més assequible.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S25: Hotels per comarca")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "03", "Oferta turística", "Planta hotelera per comarca — Places i establiments")
top10h = hotel_com[:10]
h_names  = [c[0].title()[:22] for c in top10h]
h_plazas = [c[1]['plazas'] for c in top10h]
h_estab  = [c[1]['establecimientos'] for c in top10h]
ch = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, h_names,
               [('Places', h_plazas), ('Establiments (×10)', [e*10 for e in h_estab])],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(11.0), legend=True, lpos='bottom')
ch.series[0].format.fill.fore_color.rgb = TEAL
ch.series[1].format.fill.fore_color.rgb = CYAN
tb(s, f"La comarca de {top10h[0][0].title()} concentra {fmt(top10h[0][1]['plazas'])} places hoteleres "
      f"({fmt(top10h[0][1]['establecimientos'])} establiments), el "
      f"{fmt(top10h[0][1]['plazas']/M['hoteles_plazas']*100,1)}% del total provincial. "
      f"La Safor és la principal comarca costanera amb {fmt(top10h[1][1]['plazas'])} places.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.0), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S26: VUT top municipis")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "03", "Oferta turística", "Vivendes d'Ús Turístic (VUT) — Top 12 municipis")
top12v = vut_mun[:12]
v_names  = [m[0].title()[:18] for m in top12v]
v_counts = [m[1]['count'] for m in top12v]
ch = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, v_names, [('VUT', v_counts)],
               Cm(0.5), Cm(2.3), Cm(13.5), Cm(11.0), legend=False)
ch.series[0].format.fill.fore_color.rgb = CYAN
kpi(s, Cm(14.3), Cm(2.5), Cm(8.0), Cm(2.8), fmt(M['vut_total']),
    f"VUT registrades\n{fmt(M['vut_plazas'])} places totals", bg=TEAL)
tb(s, f"{top12v[0][0].title()} ({fmt(top12v[0][1]['count'])} VUT) i "
      f"{top12v[1][0].title()} ({fmt(top12v[1][1]['count'])} VUT) lideren, "
      f"concentrades a La Safor i La Ribera. La distribució reflecteix la demanda "
      f"d'allotjament familiar en destinacions de sol i platja. "
      f"Les VUT suposen el {fmt(M['vut_plazas']/total_plazas*100,1)}% de les places reglades.",
   Cm(14.3), Cm(5.5), Cm(8.0), Cm(6.0), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S27: VUT per comarca")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "03", "Oferta turística", "Vivendes d'Ús Turístic (VUT) per comarca")
vc_names  = [c[0][:22] for c in vut_com]
vc_counts = [c[1]['count'] for c in vut_com]
ch = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, vc_names, [('VUT per comarca', vc_counts)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(11.0), legend=False)
ch.series[0].format.fill.fore_color.rgb = CYAN
tb(s, f"La Safor lidera amb {fmt(vut_com[0][1]['count'])} VUT registrades "
      f"({fmt(vut_com[0][1]['plazas'])} places), seguida de la comarca de València "
      f"({fmt(vut_com[5][1]['count'])} VUT) i El Camp de Morvedre. "
      f"Aquesta distribució contrasta amb la planta hotelera, concentrada a la capital, "
      f"evidenciant que les VUT cobreixen la demanda residencial del litoral.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.0), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S28: Càmpings i cases rurals")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "03", "Oferta turística", "Càmpings i cases rurals — distribució territorial")
top8c = camp_mun[:8]
if top8c:
    cc_n = [m[0].title()[:15] for m in top8c]
    cc_p = [m[1]['plazas'] for m in top8c]
    ch_c = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, cc_n, [('Places', cc_p)],
                     Cm(0.5), Cm(2.8), Cm(10.5), Cm(5.2), legend=False)
    ch_c.series[0].format.fill.fore_color.rgb = GREEN
    tb(s, f"CÀMPINGS: {M['campings_total']} establ. · {fmt(M['campings_plazas'])} places",
       Cm(0.5), Cm(2.4), Cm(10.5), Cm(0.6), sz=10, bold=True, color=TEAL)

top8r = rural_mun[:8]
if top8r:
    rr_n = [m[0].title()[:15] for m in top8r]
    rr_c = [m[1]['count'] for m in top8r]
    ch_r = chart_add(s, XL_CHART_TYPE.BAR_CLUSTERED, rr_n, [('Establiments', rr_c)],
                     Cm(0.5), Cm(9.0), Cm(10.5), Cm(4.5), legend=False)
    ch_r.series[0].format.fill.fore_color.rgb = TEAL
    tb(s, f"CASES RURALS: {M['casasrurales_total']} establ. · {fmt(M['casasrurales_plazas'])} places",
       Cm(0.5), Cm(8.5), Cm(10.5), Cm(0.6), sz=10, bold=True, color=TEAL)

tb_lines(s, [
    (f"Càmpings ({M['campings_total']} establiments)", True, 11),
    (f"Els {M['campings_total']} àrees de càmping sumen {fmt(M['campings_plazas'])} places. "
     f"Oliva, Sueca i Gandia concentren la major capacitat. "
     f"El litoral de La Safor és el principal nucli campista.", False, 9),
    ("", False, 6),
    (f"Cases rurals ({M['casasrurales_total']} establiments)", True, 11),
    (f"Les {M['casasrurales_total']} cases rurals ({fmt(M['casasrurales_plazas'])} places) "
     f"se distribueixen per l'interior, destacant Chulilla, Requena i Aras de los Olmos. "
     f"Complementen les destinacions de natura i turisme actiu.", False, 9),
    ("", False, 6),
    (f"Albergs: {M['albergues_total']} establ. · {fmt(M['albergues_plazas'])} places", True, 10),
],  Cm(11.5), Cm(2.5), Cm(11.0), Cm(11.0), color=DGRAY)
footer(s)

print("S29: Distribució places per tipus")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "03", "Oferta turística", "Distribució de places per tipus d'allotjament")
tipos     = ['Hotels','VUT','Càmpings','Cases rurals','Albergs']
tipo_vals = [M['hoteles_plazas'], M['vut_plazas'], M['campings_plazas'],
             M['casasrurales_plazas'], M['albergues_plazas']]
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, tipos,
               [('Places', tipo_vals)],
               Cm(1.0), Cm(2.5), Cm(14.0), Cm(11.0), legend=False)
ch.series[0].format.fill.fore_color.rgb = TEAL
rows = [f"{t}: {fmt(v)} places ({fmt(v/sum(tipo_vals)*100,1)}%)"
        for t,v in zip(tipos, tipo_vals)]
tb_lines(s, [(f"Total: {fmt(sum(tipo_vals))} places", True, 14)] + [(r, False, 11) for r in rows],
         Cm(15.5), Cm(3.0), Cm(7.0), Cm(8.0), color=DGRAY)
footer(s)

print("S30: Conclusions oferta")
s = prs.slides.add_slide(blank)
conclusions_slide(s, "03  OFERTA TURÍSTICA", [
    ("Oferta hotelera diversificada i de qualitat",
     f"Els {M['hoteles_total']} hotels ({fmt(M['hoteles_plazas'])} places) cobreixen tot l'espectre, "
     f"amb predomini de 3 i 4 estrelles. La comarca de València concentra el "
     f"{fmt(hotel_com[0][1]['plazas']/M['hoteles_plazas']*100,1)}% de la planta hotelera."),
    ("Les VUT consoliden la seva posició com a segon tipus d'allotjament",
     f"Amb {fmt(M['vut_total'])} habitatges ({fmt(M['vut_plazas'])} places), les VUT representen "
     f"el {fmt(M['vut_plazas']/total_plazas*100,1)}% de l'oferta total. Oliva i Cullera lideren."),
    ("Oferta rural i natural en expansió a l'interior",
     f"Les {M['casasrurales_total']} cases rurals i {M['campings_total']} càmpings ({fmt(M['campings_plazas'])} places) "
     f"diversifiquen l'oferta cap a productes de natura, gastronomia i turisme actiu."),
    ("Total de places supera les 113.000",
     f"Sumant tots els tipus, la província compta amb {fmt(total_plazas)} places reglades, "
     f"posicionant-la com una de les principals destinacions turístiques d'Espanya."),
], accent=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓ 4 — TREBALL TURÍSTIC
# ══════════════════════════════════════════════════════════════════════════════
print("S31: Divisor Treball")
s = slide_copy(2)
total_q3 = sum(ss[ss_dates[2]].values())
divider(s, "04", "TREBALL\nTURÍSTIC",
        f"Màxim {fmt(total_q3)} afiliats (Q3)  ·  Restauració: {fmt(rest_v[2])} afiliats (juliol)  ·  Allotjament: {fmt(aloj_v[2])}",
        color_accent=GREEN)

print("S32: KPIs empleo")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=LGRAY)
add_rect(s, Cm(0), Cm(0), W, Cm(1.8), fill=TEAL)
tb(s, "04  TREBALL TURÍSTIC · Afiliats Seguretat Social 2025", Cm(0.5), Cm(0.2), W-Cm(1), Cm(1.4),
   sz=14, bold=True, color=WHITE)
total_q1  = sum(ss[ss_dates[0]].values())
total_q4  = sum(ss[ss_dates[3]].values())
emp_kpis = [
    (fmt(total_q3),      f"Afiliats màxim Q3 (juliol)\n+{fmtp((total_q3/total_q1-1)*100)} vs Q1", TEAL),
    (fmt(rest_v[2]),     f"Restauració Q3\n{fmt(rest_v[2]/total_q3*100,1)}% del total turístic",   CYAN),
    (fmt(aloj_v[2]),     f"Allotjament Q3\npic estacional juliol",                                  GREEN),
    (fmt(total_q4),      f"Afiliats Q4 (octubre)\nbase de tardor",                                  TEAL),
    (fmt(rest_v[3]),     f"Restauració Q4\nestable respecte Q3",                                     CYAN),
    (fmt(agenc_v[0]),    f"Agències de viatge Q1\nbaixa estacionalitat",                             GREEN),
]
for i,(v,l,bg) in enumerate(emp_kpis):
    lx = Cm(0.4) + (i%3)*(bw2+gp2)
    ty = Cm(2.0) + (i//3)*(bh2+gp2)
    kpi(s, lx, ty, bw2, bh2, v, l, bg=bg)
footer(s)

print("S33: Evolució trimestral")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "04", "Treball turístic", "Evolució trimestral d'afiliats a la Seguretat Social per sector")
ch = chart_add(s, XL_CHART_TYPE.COLUMN_CLUSTERED, quarters,
               [('Allotjament', aloj_v), ('Restauració', rest_v), ('Agències viatge', agenc_v)],
               Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8), legend=True, lpos='bottom')
ch.series[0].format.fill.fore_color.rgb = TEAL
ch.series[1].format.fill.fore_color.rgb = CYAN
ch.series[2].format.fill.fore_color.rgb = GREEN
tb(s, f"La restauració és el principal ocupador ({fmt(max(rest_v))} afiliats en el pic). "
      f"L'allotjament mostra l'estacionalitat més marcada: el pic de juliol ({fmt(max(aloj_v))}) "
      f"supera en un {fmt((max(aloj_v)/min(aloj_v)-1)*100,1)}% el mínim de gener ({fmt(min(aloj_v))}). "
      f"Les agències de viatge presenten plantilla estable al voltant dels {fmt(sum(agenc_v)//4)} afiliats.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S34: Estacionalitat empleo")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=WHITE)
header(s, "04", "Treball turístic", "Estacionalitat de l'ocupació — Índex trimestral (Q1 = 100)")
aloj_idx  = [round(v/aloj_v[0]*100, 1) for v in aloj_v]
rest_idx  = [round(v/rest_v[0]*100, 1) for v in rest_v]
agenc_idx = [round(v/agenc_v[0]*100, 1) for v in agenc_v]
ch = line_chart(s, quarters,
                [('Allotjament', aloj_idx), ('Restauració', rest_idx), ('Agències viatge', agenc_idx)],
                Cm(0.5), Cm(2.3), Cm(15.5), Cm(10.8))
ch.series[0].format.line.color.rgb = TEAL
ch.series[1].format.line.color.rgb = CYAN
ch.series[2].format.line.color.rgb = GREEN
tb(s, f"L'allotjament presenta la major estacionalitat: passa de 100 a {max(aloj_idx):.0f} al Q3 "
      f"(oscil·lació del {max(aloj_idx)-100:.0f}%). La restauració mostra un rang {min(rest_idx):.0f}–{max(rest_idx):.0f}. "
      f"Les agències de viatge, amb variació de {max(agenc_idx)-min(agenc_idx):.0f} punts índex, "
      f"mantenen l'ocupació més estable al llarg de l'any.",
   Cm(16.3), Cm(2.5), Cm(6.0), Cm(6.5), sz=9.5, color=DGRAY, wrap=True)
footer(s)

print("S35: Conclusions treball")
s = prs.slides.add_slide(blank)
conclusions_slide(s, "04  TREBALL TURÍSTIC", [
    ("Sector amb capacitat per a 90.000+ afiliats en temporada alta",
     f"El pic de Q3 assoleix {fmt(total_q3)} afiliats, amb un increment estacional del "
     f"{fmtp((total_q3/total_q1-1)*100)} respecte a Q1. L'allotjament és el sector de major variació."),
    ("La restauració, base estable de l'ocupació turística",
     f"Amb {fmt(max(rest_v))} afiliats al pic i {fmt(min(rest_v))} al mínim, la restauració "
     f"representa el motor d'ocupació contínua del sector, amb baixa estacionalitat relativa."),
    ("Agències de viatge: estabilitat i especialització",
     f"Les {fmt(max(agenc_v))} agències d'afiliats mostren el perfil laboral més estable, "
     f"reflectint la naturalesa del sector: activitat distribuïda uniformement al llarg de l'any."),
    ("Potencial de desestacionalització del treball",
     f"La ràtio Q3/Q1 de {M['seasonality_ratio_ss']:.1f}% indica una estacionalitat "
     f"moderada. Polítiques de desestacionalització podrien reforçar l'ocupació de qualitat."),
], accent=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 36 — CONCLUSIONS GENERALS I PERSPECTIVA 2026
# ══════════════════════════════════════════════════════════════════════════════
print("S36: Conclusions generals")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=TEAL)
tb(s, "CONCLUSIONS GENERALS I PERSPECTIVA 2026",
   Cm(1), Cm(0.3), W-Cm(2), Cm(1.2), sz=16, bold=True, color=WHITE)

conc_gen = [
    ("01 · Rècord históric de turisme internacional",
     f"Amb {fmt(M['total_receptor_2025'])} visitants estrangers ({fmtp(M['chg_receptor'])}), la província "
     f"bat el seu màxim históric. França, Itàlia, Regne Unit, Països Baixos i Alemanya lideren. "
     f"EUA (+{fmtp(top10i[6]['chg'])}) i Polònia (+{fmtp(top10i[7]['chg'])}) emergeixen com nous mercats."),
    ("02 · Posicionament competitiu reforçat",
     f"El {fmt(M['total_receptor_2025']/M['total_turistas_2025']*100,1)}% del turisme és ja internacional, "
     f"amb estacionalitat reduïda a tardor i hivern. El diferencial nacional/internacional s'escurça."),
    ("03 · Infraestructura turística sòlida: 113.000+ places",
     f"Hotels (47k places), VUT (38k), càmpings (21k) i oferta rural (3k) garanteixen "
     f"capacitat i diversificació. L'ocupació hotelera del {M['ocupacion_hoteles_2025_avg']:.1f}% és competitiva."),
    ("04 · Treball turístic: 90.000+ afiliats en temporada alta",
     f"El sector genera ocupació estable i de qualitat al llarg de l'any. "
     f"La restauració és el gran motor, l'allotjament el de major estacionalitat."),
    ("05 · Perspectiva 2026",
     f"La tendència d'internacionalització, el creixement de les VUT i les oportunitats en mercats "
     f"emergents (EUA, Polònia) marquen les prioritats. La clau serà consolidar la demanda de "
     f"temporada baixa i mantenir la qualitat de l'oferta davant d'una competència global creixent."),
]
y = Cm(1.9)
for title, body in conc_gen:
    add_rect(s, Cm(0.8), y+Cm(0.1), Cm(0.4), Cm(0.4), fill=CYAN)
    tb(s, title, Cm(1.5), y, W-Cm(2.5), Cm(0.65), sz=11, bold=True, color=CYAN)
    tb(s, body, Cm(1.5), y+Cm(0.6), W-Cm(2.5), Cm(0.8), sz=9,
       color=RGBColor(0xDD,0xDD,0xDD), wrap=True)
    y += Cm(1.55)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 37 — GLOSARI I FONTS
# ══════════════════════════════════════════════════════════════════════════════
print("S37: Glosari")
s = prs.slides.add_slide(blank)
add_rect(s, Cm(0), Cm(0), W, H, fill=LGRAY)
add_rect(s, Cm(0), Cm(0), W, Cm(1.8), fill=TEAL)
tb(s, "GLOSARI DE TERMES I FONTS DE DADES", Cm(0.5), Cm(0.2), W-Cm(1), Cm(1.4),
   sz=14, bold=True, color=WHITE)
glos = [
    ("Turisme receptor", "Visitants no residents a Espanya que viatgen a la província des de l'estranger (INE – Enquesta de Moviments Turístics)."),
    ("Turisme intern",    "Visitants residents a Espanya que viatgen dins del territori nacional cap a la província de València."),
    ("Pernoctació",       "Cada nit que un viatger pernocta en un establiment d'allotjament turístic reglat (font: EOH–INE)."),
    ("VUT",               "Vivenda d'Ús Turístic: immoble cedit amb finalitats turístiques inscrit al Registre de Turisme de la C. Valenciana."),
    ("Grau d'ocupació",   "Percentatge de places o habitacions ocupades respecte al total disponible estimat (EOH–INE)."),
    ("Estada mitjana",    "Nombre de nits de mitjana que un viatger pernocta a la destinació per visita."),
]
fonts = [
    ("INE – Estadística Moviments Turístics", "Estadística mensual de viajants i pernoctacions per CCAA i província. Font principal d'afluència, mercats i ocupació."),
    ("GVA Turisme – Registre VIVTUR",          "Registre oficial d'allotjaments de la C. Valenciana: hotels, VUT, càmpings, cases rurals i albergs."),
    ("Seguretat Social – Afiliats",            "Nombre d'afiliats per sector d'activitat (allotjament, restauració, agències). Freqüència trimestral."),
    ("Smart Office Diputació de València",     "Plataforma d'intel·ligència turística provincial que integra totes les fonts per a l'anàlisi agregada."),
    ("Períodes analitzats",                   "ENE–DIC 2025. Comparatives respecte a ENE–DIC 2024. Oferta: dades registrals novembre 2025."),
    ("Nota metodológica",                     "Turisme receptor i intern: dades de moviments a nivell municipal agregades a nivell provincial i comarcal."),
]
y_g = Cm(2.1)
for i, ((t1,d1),(t2,d2)) in enumerate(zip(glos, fonts)):
    tb(s, t1, Cm(0.5), y_g, Cm(11.0), Cm(0.65), sz=10, bold=True, color=TEAL)
    tb(s, d1, Cm(0.5), y_g+Cm(0.6), Cm(11.0), Cm(0.85), sz=8.5, color=DGRAY, wrap=True)
    tb(s, t2, Cm(12.0), y_g, Cm(11.0), Cm(0.65), sz=10, bold=True, color=TEAL)
    tb(s, d2, Cm(12.0), y_g+Cm(0.6), Cm(11.0), Cm(0.85), sz=8.5, color=DGRAY, wrap=True)
    y_g += Cm(1.35)
footer(s)

# ══════════════════════════════════════════════════════════════════════════════
# CONTRAPORTADA
# ══════════════════════════════════════════════════════════════════════════════
print("S38: Contraportada")
s = slide_copy(8)

# ══════════════════════════════════════════════════════════════════════════════
# ELIMINAR SLIDES DE PLANTILLA ORIGINALS (primers N_TPL)
# ══════════════════════════════════════════════════════════════════════════════
print("\nEliminant slides de plantilla originals...")
sld_id_lst = prs.slides._sldIdLst
for _ in range(N_TPL):
    sld_id_lst.remove(sld_id_lst[0])

# ══════════════════════════════════════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════════════════════════════════════
print(f"\nGuardant {OUTPUT}...")
prs.save(OUTPUT)
print(f"\n✓ Informe DIVAL 2025 v2 generat: {len(prs.slides)} slides")
print(f"  Arxiu: {OUTPUT}")
